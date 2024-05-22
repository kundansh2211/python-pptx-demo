from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Add a slide with a title and chart layout
slide_layout = presentation.slide_layouts[5]  # Assuming layout 5 has title and content
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Peanut Production in 5 States of USA (2020-2022)"

# Define chart data
chart_data = CategoryChartData()
chart_data.categories = ['Georgia', 'Illinois', 'Massachusetts', 'New York', 'Texas']
chart_data.add_series('2020', (19.2, 21.4, 16.7, 23.5, 18.3))
chart_data.add_series('2021', (22.4, 25.3, 18.9, 27.1, 20.8))
chart_data.add_series('2022', (24.5, 27.8, 21.1, 29.6, 22.7))

# Add chart to the slide
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Customize chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# chart.value_axis.has_major_gridlines = True
value_axis = chart.value_axis
value_axis.maximum_scale = 30.0
value_axis.minimum_scale = 0.0
value_axis.major_unit = 5.0

category_axis = chart.category_axis
# category_axis.has_major_gridlines = True

# Save the presentation
presentation.save('PPT/peanut_production1.pptx')

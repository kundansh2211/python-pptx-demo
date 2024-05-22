from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Create a presentation object
prs = Presentation()

# Add a title slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout index for title slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sample Presentation"
subtitle.text = "This is a subtitle"

# Add a slide with text
slide_layout = prs.slide_layouts[1]  # 1 is the layout index for title and content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Slide with Text"
content.text = "This is a sample text content."

# Add a slide with an image
slide_layout = prs.slide_layouts[5]  # 5 is the layout index for a blank slide
slide = prs.slides.add_slide(slide_layout)
img_path = 'images/bar.jpg'
left = Inches(1)
top = Inches(1)
slide.shapes.add_picture(img_path, left, top, width=Inches(4), height=Inches(3))

# Add a slide with a chart
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
chart_data = CategoryChartData()
chart_data.categories = ['Category 1', 'Category 2', 'Category 3']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Customize chart 
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.plots[0].has_data_labels = True
chart.plots[0].data_labels.font.size = Pt(12)

# Save the presentation
prs.save('PPT/sample_presentation.pptx')

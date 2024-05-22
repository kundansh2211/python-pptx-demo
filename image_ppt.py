import os
from pptx import Presentation


# Define the directory where images are stored
image_dir = 'images'

# List of images to be converted to ppt
image_filenames = ['bar.jpg', 'images.png', 'peanuts.jpg']

# Create a presentation object
presentation = Presentation()

# Loop through the image filenames and add slides with images
for filename in image_filenames:
    # Construct the full path to the image file
    image_path = os.path.join(image_dir, filename)

    # Check if the file exists
    if not os.path.isfile(image_path):
        print(f"File not found: {image_path}")
        continue

    try:
        # Create a slide with blank layout
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        # Define the height and width of the image
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        image_width = slide_width
        image_height = slide_height

        # Calculate the position to center the image on the slide
        left = (slide_width - image_width) / 2
        top = (slide_height - image_height) / 2

        # Add the image to the slide
        slide.shapes.add_picture(
            image_path, left, top, image_width, image_height)
    except Exception as e:
        print(f"Error adding image {image_path}: {e}")

# Save the presentation
presentation.save('PPT/image_ppt.pptx')
